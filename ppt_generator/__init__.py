"""
PPT生成器核心模块

借鉴 NotebookLM 和 Nano Banana Pro 最佳实践：
- 两阶段大纲生成（文档分析 + 大纲生成）
- 结构化 Prompt 模板
- 智能错误处理与降级
- 缓存机制（大纲 + 图片）
- 批量图片生成与风格锚定
"""

import asyncio
import os
import json
import logging
from datetime import datetime
from typing import List, Dict, Optional

from .outline_generator import OutlineGenerator
from .claude_client import ClaudeClient, OpenaiClient, UnifiedAIClient
from .cache_manager import CacheManager
from .batch_generator import BatchImageGenerator
from .error_handler import SmartErrorHandler
from .template_loader import get_template_presets

# 优先使用内部版本，如果不存在则使用官方版本
try:
    from .slide_generator import SlideGenerator, ImageGenerationTool
except ImportError:
    from .slide_generator_official import SlideGenerator, ImageGenerationTool

__version__ = "2.0.0"

logger = logging.getLogger(__name__)


class PPTGenerator:
    """
    PPT生成器主类

    核心功能：
    - 两阶段大纲生成（文档分析 + 大纲生成）
    - 风格锚定批量图片生成
    - 缓存机制（大纲 + 图片）
    - 智能错误处理与降级
    - 模板预设支持
    """

    def __init__(
        self,
        api_key: str = None,
        provider: str = "Claude",
        base_url: str = None,
        enable_cache: bool = True,
        cache_dir: str = ".cache/ppt_generator"
    ):
        """
        初始化PPT生成器

        Args:
            api_key: API密钥，如果为None则从环境变量读取
            provider: AI提供商 (Claude/Openai/other)
            base_url: 自定义API基础URL
            enable_cache: 是否启用缓存
            cache_dir: 缓存目录
        """
        # 初始化 LLM 客户端
        if provider == "Claude":
            self.llm_client = ClaudeClient(api_key)
        elif provider == "Openai":
            self.llm_client = OpenaiClient(api_key=api_key, base_url=base_url)
        else:
            self.llm_client = UnifiedAIClient(api_key=api_key, base_url=base_url)

        # 核心组件
        self.outline_generator = OutlineGenerator(self.llm_client)
        self.image_tool = ImageGenerationTool()
        self.slide_generator = SlideGenerator(self.image_tool)

        # 增强组件
        self.cache_manager = CacheManager(cache_dir) if enable_cache else None
        self.batch_generator = BatchImageGenerator(self.image_tool, self.cache_manager)
        self.error_handler = SmartErrorHandler()

        self.enable_cache = enable_cache

    def generate_ppt(
        self,
        reference_text: str,
        style_requirements: str,
        output_dir: str = "output",
        model: str = "deepseek-chat",
        audience_profile: Dict = None,
        brand_guidelines: Dict = None,
        brand_references: List[str] = None,
        use_cache: bool = True,
        template_preset: str = None
    ) -> dict:
        """
        同步生成PPT

        借鉴 NotebookLM + Nano Banana Pro 最佳实践：
        - 两阶段大纲生成（文档分析 → 大纲生成）
        - 风格锚定批量图片生成
        - 缓存机制（大纲 + 图片）
        - 智能错误处理与降级

        Args:
            reference_text: 参考文本
            style_requirements: 风格要求
            output_dir: 输出目录
            model: 使用的模型
            audience_profile: 目标受众信息
            brand_guidelines: 品牌规范
            brand_references: 品牌参考图路径列表（最多14张）
            use_cache: 是否使用缓存
            template_preset: 模板预设名称，可选值:
                - "business_pitch": 商业路演
                - "technical_report": 技术报告
                - "product_launch": 产品发布
                - "training": 培训课程
                - "quarterly_review": 季度汇报
                - "project_proposal": 项目提案
                - "company_intro": 公司介绍
                - "academic": 学术演讲

        Returns:
            dict: 生成结果，包含文件路径和生成信息
        """
        return asyncio.run(self.generate_ppt_async(
            reference_text, style_requirements, output_dir, model,
            audience_profile, brand_guidelines, brand_references, use_cache,
            template_preset=template_preset
        ))

    async def generate_ppt_async(
        self,
        reference_text: str,
        style_requirements: str,
        output_dir: str = "output",
        model: str = "deepseek-chat",
        audience_profile: Dict = None,
        brand_guidelines: Dict = None,
        brand_references: List[str] = None,
        use_cache: bool = True,
        max_concurrent: int = 4,
        template_preset: str = None
    ) -> dict:
        """
        异步生成PPT（推荐使用）

        借鉴 NotebookLM + Nano Banana Pro 最佳实践：
        1. 两阶段大纲生成（文档分析 → 大纲生成）
        2. 缓存机制（大纲 + 图片）
        3. 风格锚定批量生成
        4. 智能错误处理与降级

        Args:
            reference_text: 参考文本
            style_requirements: 风格要求
            output_dir: 输出目录
            model: 使用的模型
            audience_profile: 目标受众信息
            brand_guidelines: 品牌规范
            brand_references: 品牌参考图路径列表
            use_cache: 是否使用缓存
            max_concurrent: 最大并发数
            template_preset: 模板预设名称

        Returns:
            dict: 生成结果，包含文件路径和详细信息
        """
        generation_info = {
            "mode": "enhanced",
            "two_stage": True,
            "cache_used": False,
            "style_anchored": True,
            "template_preset": template_preset
        }

        # ===== 阶段1：检查缓存 =====
        outline_result = None
        if use_cache and self.cache_manager:
            logger.info("检查大纲缓存...")
            outline_result = self.cache_manager.get_cached_outline(
                reference_text, style_requirements, model
            )
            if outline_result:
                logger.info(f"命中大纲缓存，共 {len(outline_result.get('slides', []))} 页")
                generation_info["cache_used"] = True

        # ===== 阶段2：两阶段大纲生成 =====
        if not outline_result:
            logger.info("开始两阶段大纲生成...")
            try:
                outline_result = self.outline_generator.generate_outline_two_stage(
                    reference_text=reference_text,
                    style_requirements=style_requirements,
                    audience_profile=audience_profile,
                    brand_guidelines=brand_guidelines,
                    model=model,
                    template_preset=template_preset
                )
                logger.info(f"两阶段大纲生成完成，共 {len(outline_result.get('slides', []))} 页")

                # 缓存大纲
                if use_cache and self.cache_manager:
                    self.cache_manager.cache_outline(
                        reference_text, style_requirements, model, outline_result
                    )
                    logger.info("大纲已缓存")

            except Exception as e:
                logger.warning(f"两阶段生成失败，降级到单阶段: {e}")
                outline_result = self.outline_generator.generate_outline(
                    reference_text, style_requirements, model=model, template_preset=template_preset
                )
                generation_info["two_stage"] = False

        # ===== 阶段3：批量图片生成（带风格锚定）=====
        logger.info("开始批量图片生成（风格锚定模式）...")

        # 获取模板预设的风格提示
        style_hints = None
        template_presets = get_template_presets()
        if template_preset and template_preset in template_presets:
            preset = template_presets[template_preset]
            style_hints = preset.get('style_hints')
            if style_hints:
                logger.info(f"使用预设风格: {preset.get('name', template_preset)} - {style_hints.get('background', '')}")

        try:
            slides = outline_result.get('slides', [])
            batch_results = await self.batch_generator.generate_with_style_consistency(
                slides=slides,
                outline_result=outline_result,
                style_requirements=style_requirements,
                output_dir=output_dir,
                brand_references=brand_references,
                max_concurrent=max_concurrent,
                style_hints=style_hints
            )

            # 转换为标准格式
            slides_data = []
            for result in batch_results:
                slides_data.append({
                    'success': result.success,
                    'file_path': result.file_path,
                    'error': result.error,
                    'title': slides[result.slide_index].get('title', '') if result.slide_index < len(slides) else '',
                    'from_cache': result.from_cache,
                    'retries': result.retries
                })

        except Exception as e:
            logger.warning(f"批量生成失败，降级到标准模式: {e}")
            generation_info["style_anchored"] = False
            slides_data = await self._generate_ppt_slides_fallback(
                outline_result, style_requirements, output_dir, max_concurrent,
                style_hints=style_hints  # 传递 style_hints 到降级模式
            )

        # ===== 阶段4：保存结果 =====
        result = self._save_ppt(slides_data, outline_result, output_dir)
        result["generation_info"] = generation_info

        # 统计缓存命中
        cache_hits = sum(1 for s in slides_data if s.get('from_cache', False))
        if cache_hits > 0:
            result["cache_hits"] = cache_hits
            logger.info(f"图片缓存命中: {cache_hits}/{len(slides_data)}")

        return result

    async def _generate_ppt_slides_fallback(
        self,
        outline_result: dict,
        style_requirements: str,
        output_dir: str,
        max_concurrent: int = 4,
        timeout_per_slide: float = 120.0,
        style_hints: dict = None
    ) -> List[dict]:
        """
        降级模式：并发生成PPT幻灯片（无风格锚定）

        Args:
            outline_result: 大纲结果字典
            style_requirements: 风格要求
            output_dir: 输出目录
            max_concurrent: 最大并发数
            timeout_per_slide: 单个幻灯片超时时间（秒）
            style_hints: 模板预设的风格提示（可选）

        Returns:
            List[dict]: 幻灯片数据列表
        """
        slides = outline_result['slides']
        total_slides = len(slides)

        logger.info(f"降级模式：并发生成PPT页面，共{total_slides}页，并发数:{max_concurrent}")
        if style_hints:
            logger.info(f"使用预设风格提示: {style_hints.get('background', '')[:30]}...")

        semaphore = asyncio.Semaphore(max_concurrent)

        async def generate_single_slide(i: int, slide_info: dict) -> tuple:
            async with semaphore:
                try:
                    result = await asyncio.wait_for(
                        self.slide_generator.generate_slide_as_image(
                            slide_info, i, outline_result, style_requirements, output_dir,
                            style_hints=style_hints  # 传递 style_hints
                        ),
                        timeout=timeout_per_slide
                    )

                    if result.get("success"):
                        logger.info(f"✓ 第 {i+1}/{total_slides} 页生成成功 - {slide_info.get('title', 'Untitled')}")
                        return (i, {'success': True, 'file_path': result.get('file_path'), 'title': slide_info.get('title', '')})
                    else:
                        error_msg = result.get('error', 'Unknown error')
                        logger.warning(f"✗ 第 {i+1}/{total_slides} 页生成失败: {error_msg}")
                        return (i, {'success': False, 'error': error_msg, 'title': slide_info.get('title', '')})

                except asyncio.TimeoutError:
                    logger.error(f"✗ 第 {i+1}/{total_slides} 页生成超时")
                    return (i, {'success': False, 'error': '生成超时', 'title': slide_info.get('title', '')})

                except Exception as e:
                    logger.exception(f"✗ 第 {i+1}/{total_slides} 页生成异常: {str(e)}")
                    return (i, {'success': False, 'error': str(e), 'title': slide_info.get('title', '')})

        tasks = [generate_single_slide(i, slide) for i, slide in enumerate(slides)]
        results = await asyncio.gather(*tasks)

        sorted_results = sorted(results, key=lambda x: x[0])
        slides_data = [data for _, data in sorted_results]

        logger.info(f"✓ 所有 {total_slides} 页幻灯片生成完成！")
        return slides_data

    def _save_ppt(self, slides_data: list, outline: dict, output_dir: str) -> dict:
        """
        保存PPT为.pptx文件

        Args:
            slides_data: 幻灯片数据列表 (包含图片路径或错误信息)
            outline: 大纲信息
            output_dir: 输出目录

        Returns:
            dict: 保存结果信息
                - pptx_file: .pptx文件路径
                - outline_file: 大纲JSON路径
                - total_slides: 总页数
                - success_slides: 成功页数
                - error_slides: 失败页面列表
                - timestamp: 时间戳
        """
        from pptx import Presentation
        from pptx.util import Inches

        os.makedirs(output_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # 创建PPT对象 (16:9 比例)
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # 16:9 标准宽度
        prs.slide_height = Inches(7.5)     # 16:9 标准高度

        success_count = 0
        error_slides = []

        for i, slide_data in enumerate(slides_data):
            # 添加空白幻灯片
            blank_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(blank_layout)

            if slide_data.get('success') and slide_data.get('file_path'):
                image_path = slide_data['file_path']
                if os.path.exists(image_path):
                    # 成功的幻灯片：插入图片作为全屏背景
                    slide.shapes.add_picture(
                        image_path,
                        Inches(0), Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    success_count += 1
                else:
                    # 图片文件不存在
                    self._add_error_text_to_slide(
                        slide, prs, i + 1,
                        slide_data.get('title', ''),
                        f"图片文件不存在: {image_path}"
                    )
                    error_slides.append({
                        'page': i + 1,
                        'title': slide_data.get('title', ''),
                        'error': f"图片文件不存在: {image_path}"
                    })
            else:
                # 失败的幻灯片：添加错误文本
                error_msg = slide_data.get('error', '未知错误')
                self._add_error_text_to_slide(
                    slide, prs, i + 1,
                    slide_data.get('title', ''),
                    error_msg
                )
                error_slides.append({
                    'page': i + 1,
                    'title': slide_data.get('title', ''),
                    'error': error_msg
                })

        # 保存PPT文件
        pptx_filename = f"presentation_{timestamp}.pptx"
        pptx_path = os.path.join(output_dir, pptx_filename)
        prs.save(pptx_path)

        # 保存大纲信息
        outline_file = os.path.join(output_dir, f"outline_{timestamp}.json")
        with open(outline_file, 'w', encoding='utf-8') as f:
            json.dump(outline, f, ensure_ascii=False, indent=2)

        logger.info(f"PPT已保存: {pptx_path}，成功 {success_count}/{len(slides_data)} 页")

        return {
            'pptx_file': pptx_path,
            'outline_file': outline_file,
            'total_slides': len(slides_data),
            'success_slides': success_count,
            'error_slides': error_slides,
            'timestamp': timestamp
        }

    def _add_error_text_to_slide(
        self,
        slide,
        prs,
        page_number: int,
        title: str,
        error_msg: str
    ):
        """
        在幻灯片上添加错误提示文本

        Args:
            slide: 幻灯片对象
            prs: Presentation对象
            page_number: 页码
            title: 幻灯片标题
            error_msg: 错误信息
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # 添加红色背景形状
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 53, 69)  # Bootstrap danger red
        shape.line.fill.background()

        # 添加标题文本框
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "页面生成错误"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # 添加页面信息
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8), Inches(11.333), Inches(0.8)
        )
        tf2 = info_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"第 {page_number} 页 - {title or '无标题'}"
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER

        # 添加错误详情
        detail_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(11.333), Inches(1)
        )
        tf3 = detail_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = f"错误信息：{error_msg}"
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(255, 200, 200)
        p3.alignment = PP_ALIGN.CENTER

    def clear_cache(self, older_than_days: int = None):
        """
        清除缓存

        Args:
            older_than_days: 清除多少天前的缓存，None表示全部清除
        """
        if self.cache_manager:
            if older_than_days:
                self.cache_manager.cleanup_expired()
                logger.info(f"已清除过期缓存（{self.cache_manager.cache_ttl.days}天前）")
            else:
                import shutil
                if self.cache_manager.cache_dir.exists():
                    shutil.rmtree(self.cache_manager.cache_dir)
                    logger.info("已清除所有缓存")
        else:
            logger.warning("缓存未启用")

    def get_cache_stats(self) -> Optional[Dict]:
        """
        获取缓存统计信息

        Returns:
            Dict: 缓存统计
                - outline_count: 大纲缓存数量
                - image_count: 图片缓存数量
                - total_size_mb: 总大小(MB)
        """
        if not self.cache_manager:
            return None

        stats = {
            "outline_count": 0,
            "image_count": 0,
            "total_size_mb": 0.0
        }

        try:
            # 统计大纲缓存
            outline_dir = self.cache_manager.outline_cache_dir
            if outline_dir.exists():
                outline_files = list(outline_dir.glob("*.json"))
                stats["outline_count"] = len(outline_files)
                for f in outline_files:
                    stats["total_size_mb"] += f.stat().st_size / (1024 * 1024)

            # 统计图片缓存
            image_dir = self.cache_manager.image_cache_dir
            if image_dir.exists():
                image_files = list(image_dir.glob("*"))
                stats["image_count"] = len(image_files)
                for f in image_files:
                    if f.is_file():
                        stats["total_size_mb"] += f.stat().st_size / (1024 * 1024)

            stats["total_size_mb"] = round(stats["total_size_mb"], 2)

        except Exception as e:
            logger.warning(f"获取缓存统计失败: {e}")

        return stats

    @staticmethod
    def list_template_presets() -> List[Dict]:
        """
        列出所有可用的模板预设

        Returns:
            List[Dict]: 预设列表，每项包含:
                - key: 预设键名（用于 template_preset 参数）
                - name: 预设中文名称
                - description: 预设描述
        """
        return [
            {
                "key": key,
                "name": preset.get("name", key),
                "description": preset.get("description", "")
            }
            for key, preset in get_template_presets().items()
        ]

    @staticmethod
    def get_template_preset_info(preset_name: str) -> Optional[Dict]:
        """
        获取指定预设的详细信息

        Args:
            preset_name: 预设名称

        Returns:
            Dict: 预设详细信息，包含 name, description, sequence, narrative
        """
        return get_template_presets().get(preset_name)
